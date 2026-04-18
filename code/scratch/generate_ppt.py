"""
PPT Generator — Delhi AQ Forecasting Project
DMS673: Applied Machine Learning | IIT Kanpur
Author: Guna Venkat Doddi (251140009)

Run: python scratch/generate_ppt.py
Requires: pip install python-pptx Pillow
Output: delhi_aq_presentation.pptx
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as pns
from lxml import etree

# ─── Paths ───────────────────────────────────────────────────────────────────
BASE = Path(__file__).parent.parent
PLOTS = BASE / "code" / "plots"
OUT   = BASE / "delhi_aq_presentation.pptx"

# ─── Colour Palette ──────────────────────────────────────────────────────────
C_BG        = RGBColor(0x0D, 0x12, 0x1F)   # near-black navy
C_SURFACE   = RGBColor(0x16, 0x1E, 0x35)   # slightly lighter navy
C_ACCENT    = RGBColor(0xFF, 0x8C, 0x00)   # deep amber
C_ACCENT2   = RGBColor(0xFF, 0xC0, 0x40)   # light amber
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT     = RGBColor(0xCC, 0xD6, 0xF1)   # light lavender-white
C_RED       = RGBColor(0xFF, 0x4D, 0x4D)
C_GREEN     = RGBColor(0x4C, 0xD9, 0x7F)
C_YELLOW    = RGBColor(0xFF, 0xD7, 0x00)
C_GRAY      = RGBColor(0x88, 0x99, 0xBB)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank


# ════════════════════════════════════════════════════════════════════════════
# HELPER FUNCTIONS
# ════════════════════════════════════════════════════════════════════════════

def add_slide():
    slide = prs.slides.add_slide(blank_layout)
    return slide

def bg(slide, color=C_BG):
    """Fill slide background with solid colour."""
    bg_el = slide.background
    fill  = bg_el.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill_color=C_SURFACE, alpha=None):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE,  # not used, just placeholder
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

from pptx.enum.shapes import MSO_SHAPE_TYPE
import pptx.shapes.autoshape

def add_rect(slide, l, t, w, h, fill=C_SURFACE, line=None):
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        if False else 1,  # 1 = rectangle
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             size=20, bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txBox

def add_multiline(slide, lines, l, t, w, h, size=16, color=C_LIGHT, spacing=1.15):
    """lines: list of (text, bold, color) tuples or plain strings."""
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            text, bold, col = item, False, color
        else:
            text, bold, col = item[0], item[1] if len(item)>1 else False, item[2] if len(item)>2 else color
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.color.rgb = col
    return txBox

def add_image(slide, path, l, t, w, h=None):
    """Add image; if path doesn't exist, add placeholder rect instead."""
    p = PLOTS / path if not os.path.isabs(str(path)) else Path(path)
    if p.exists():
        if h:
            pic = slide.shapes.add_picture(str(p), Inches(l), Inches(t), Inches(w), Inches(h))
        else:
            pic = slide.shapes.add_picture(str(p), Inches(l), Inches(t), width=Inches(w))
        return pic
    else:
        # placeholder
        box = add_rect(slide, l, t, w, h or 2.5, fill=RGBColor(0x22, 0x2E, 0x4A))
        add_text(slide, f"[{path}]", l+0.1, t+0.1, w-0.2, (h or 2.5)-0.2,
                 size=9, color=C_GRAY, align=PP_ALIGN.CENTER)
        return box

def slide_header(slide, title, subtitle=None, accent_bar=True):
    """Standard slide header: accent bar on left + title."""
    if accent_bar:
        add_rect(slide, 0, 0, 0.08, 7.5, fill=C_ACCENT)
    add_text(slide, title, 0.25, 0.18, 12.7, 0.65,
             size=28, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.25, 0.82, 12.7, 0.4,
                 size=14, bold=False, color=C_ACCENT2, align=PP_ALIGN.LEFT)
    # thin divider line
    add_rect(slide, 0.25, 1.25, 13.0, 0.025, fill=C_ACCENT)

def add_bullet_box(slide, items, l, t, w, h,
                   bullet="▸", size=15, color=C_LIGHT, title=None, title_color=C_ACCENT2):
    """Render a list of bullet items in a text box."""
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(size + 1)
        run.font.bold = True
        run.font.color.rgb = title_color
        first = False
    for item in items:
        if isinstance(item, str):
            text, bold, col = item, False, color
        else:
            text = item[0]
            bold = item[1] if len(item) > 1 else False
            col  = item[2] if len(item) > 2 else color
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"{bullet}  {text}"
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.color.rgb = col
    return txBox

def add_table_simple(slide, headers, rows, l, t, w, h,
                     hdr_fill=C_ACCENT, row_fill=C_SURFACE,
                     alt_fill=RGBColor(0x1A, 0x26, 0x42),
                     text_size=12):
    cols = len(headers)
    rws  = len(rows) + 1
    tbl  = slide.shapes.add_table(rws, cols, Inches(l), Inches(t),
                                   Inches(w), Inches(h)).table
    col_w = w / cols
    for i, c in enumerate(tbl.columns):
        c.width = Inches(col_w)

    def _cell(cell, txt, fill, bold=False, color=C_WHITE, align=PP_ALIGN.CENTER):
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
        tf = cell.text_frame
        tf.word_wrap = True
        p  = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(txt)
        run.font.size  = Pt(text_size)
        run.font.bold  = bold
        run.font.color.rgb = color

    for ci, h_txt in enumerate(headers):
        _cell(tbl.cell(0, ci), h_txt, hdr_fill, bold=True)

    for ri, row in enumerate(rows):
        fill = row_fill if ri % 2 == 0 else alt_fill
        for ci, val in enumerate(row):
            star = isinstance(val, str) and "⭐" in val
            fail = isinstance(val, str) and ("−0.88" in val or "−0.71" in val)
            col  = C_ACCENT2 if star else (C_RED if fail else C_WHITE)
            _cell(tbl.cell(ri+1, ci), val, fill, bold=star, color=col)

def add_metric_card(slide, l, t, w, h, value, label, value_color=C_ACCENT2):
    add_rect(slide, l, t, w, h, fill=C_SURFACE,
             line=RGBColor(0x3A, 0x4E, 0x78))
    add_text(slide, value, l+0.1, t+0.15, w-0.2, h*0.55,
             size=22, bold=True, color=value_color, align=PP_ALIGN.CENTER)
    add_text(slide, label, l+0.05, t+h*0.58, w-0.1, h*0.38,
             size=10, color=C_GRAY, align=PP_ALIGN.CENTER)

def slide_number(slide, num, total=19):
    add_text(slide, f"{num} / {total}", 12.5, 7.15, 0.75, 0.3,
             size=9, color=C_GRAY, align=PP_ALIGN.RIGHT)

def footer_bar(slide, text="DMS673 · Applied Machine Learning · IIT Kanpur · Guna Venkat Doddi"):
    add_rect(slide, 0, 7.25, 13.33, 0.25, fill=RGBColor(0x08, 0x0D, 0x18))
    add_text(slide, text, 0.3, 7.27, 12.5, 0.22,
             size=8, color=C_GRAY, align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)

# accent bar left
add_rect(s, 0, 0, 0.5, 7.5, fill=C_ACCENT)

# background image (faded monthly boxplot)
img_path = PLOTS / "03_monthly_boxplots_all_stations.png"
if img_path.exists():
    pic = s.shapes.add_picture(str(img_path), Inches(5.5), Inches(0),
                                Inches(7.83), Inches(7.5))
    # apply transparency via XML
    sp_tree = s.shapes._spTree
    sp_tree.remove(pic._element)
    sp_tree.insert(2, pic._element)

# dark overlay on right side
add_rect(s, 5.5, 0, 7.83, 7.5, fill=RGBColor(0x0D, 0x12, 0x1F))

# title block
add_text(s, "Spatiotemporal Air Quality", 0.7, 1.0, 11.5, 0.9,
         size=38, bold=True, color=C_WHITE)
add_text(s, "Forecasting in Delhi", 0.7, 1.85, 11.5, 0.85,
         size=38, bold=True, color=C_ACCENT)
add_text(s, "A Multi-Model Approach with Uncertainty Estimation & Explainability",
         0.7, 2.85, 9.5, 0.6, size=17, color=C_LIGHT)

# divider
add_rect(s, 0.7, 3.6, 5.5, 0.04, fill=C_ACCENT)

# meta info
add_text(s, "DMS673: Applied Machine Learning", 0.7, 3.78, 8, 0.4,
         size=14, color=C_ACCENT2)
add_text(s, "Indian Institute of Technology Kanpur", 0.7, 4.18, 8, 0.4,
         size=14, color=C_LIGHT)
add_text(s, "Guna Venkat Doddi  |  Roll No. 251140009  |  April 2026",
         0.7, 4.58, 8, 0.4, size=13, color=C_GRAY)

# station map inset
add_image(s, "07_spatial_adjacency.png", 9.0, 3.8, 3.9, 3.2)
add_text(s, "9 CPCB Monitoring Stations · Delhi NCT", 9.0, 6.95, 3.9, 0.28,
         size=8, color=C_GRAY, align=PP_ALIGN.CENTER)

slide_number(s, 1)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM & MOTIVATION
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
slide_header(s, "The Problem & Motivation",
             "Why do we need forecasts, not just readings?")

# left bullets
add_bullet_box(s, [
    ("Delhi ranks among world's most polluted cities", True, C_ACCENT2),
    "PM2.5 (<2.5 µm) penetrates deep into the respiratory system",
    "Current systems provide readings — not forecasts or what-if insights",
    "",
    ("Three Core Objectives:", True, C_ACCENT2),
    ("🔮  Forecast  —  Predict PM2.5/AQI: 1 hour → 1 week ahead", False, C_WHITE),
    ("🌬️  Simulate  —  What-If engine: 'What if wind doubles?'", False, C_WHITE),
    ("🏙️  Govern   —  Anomaly detection & station clustering", False, C_WHITE),
], 0.3, 1.35, 5.8, 5.7, bullet="", size=15)

# right plot
add_image(s, "08_aqi_category_by_year.png", 6.3, 1.38, 4.5, 3.5)
add_text(s, "AQI category distribution 2021–2025", 6.3, 4.85, 4.5, 0.3,
         size=9, color=C_GRAY, align=PP_ALIGN.CENTER)

# station ranking inset
add_image(s, "03_station_ranking_pm25.png", 11.0, 1.38, 2.0, 3.5)
add_text(s, "Station ranking", 11.0, 4.85, 2.0, 0.3,
         size=9, color=C_GRAY, align=PP_ALIGN.CENTER)

footer_bar(s); slide_number(s, 2)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — DATASET OVERVIEW
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
slide_header(s, "Dataset Overview", "CPCB multi-station air quality data · 2021–2025")

# metric cards
cards = [
    ("9", "Monitoring Stations"),
    ("4 yrs", "Data Span (2021–2025)"),
    ("393,735", "Hourly Records"),
    ("16,096", "Daily Records"),
    ("15", "Variables / Station"),
    ("15-min", "Raw Sampling Rate"),
]
cw, ch = 2.0, 1.2
for i, (val, lbl) in enumerate(cards):
    col = i % 3
    row = i // 3
    add_metric_card(s, 0.3 + col*2.1, 1.4 + row*1.35, cw, ch, val, lbl)

# station list
add_text(s, "Stations: Anand Vihar · Ashok Vihar · Bawana · Dwarka-Sec 8 · Jahangirpuri · Mundka · Punjabi Bagh · Rohini · Wazirpur",
         0.3, 4.05, 6.35, 0.55, size=11, color=C_LIGHT)

# rolling mean plot
add_image(s, "04_rolling_mean_citywide.png", 0.3, 4.65, 6.35, 2.35)
add_text(s, "4-year city-wide PM2.5 rolling trend", 0.3, 6.95, 6.35, 0.28,
         size=9, color=C_GRAY, align=PP_ALIGN.CENTER)

# station map
add_image(s, "07_spatial_adjacency.png", 6.85, 1.38, 6.2, 5.7)
add_text(s, "Station network — spatial adjacency", 6.85, 7.05, 6.2, 0.28,
         size=9, color=C_GRAY, align=PP_ALIGN.CENTER)

footer_bar(s); slide_number(s, 3)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — PREPROCESSING PIPELINE
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
slide_header(s, "Data Pipeline & Preprocessing",
             "From raw 15-min CPCB CSVs to a clean modelling-ready matrix")

# pipeline steps (left column)
steps = [
    ("① Frequency Alignment", "15-min → Hourly & Daily via mean aggregation"),
    ("② Timestamp Correction", "Fixed inconsistent datetime formats in 2023–2024 data"),
    ("③ KNN Imputation", "Missing values (sensor downtime) filled with inverse-distance KNN\n   x̂ᵢ = Σ wⱼxⱼ / Σ wⱼ,   wⱼ = 1/dist(xᵢ,xⱼ)"),
    ("④ Outlier Treatment", "99th-percentile clip per-station\n   Diwali spikes >800 µg/m³ preserved as signal"),
]
y = 1.4
for title, desc in steps:
    add_text(s, title, 0.3, y, 5.9, 0.38, size=14, bold=True, color=C_ACCENT2)
    add_text(s, desc,  0.3, y+0.35, 5.9, 0.72, size=12, color=C_LIGHT)
    add_rect(s, 0.3, y+1.08, 5.9, 0.025, fill=RGBColor(0x2A, 0x3B, 0x60))
    y += 1.22

# before/after outlier plot (right)
add_text(s, "Before & After — 99th-pct Clip (Mundka)", 6.4, 1.38, 6.7, 0.38,
         size=13, bold=True, color=C_ACCENT2)
add_image(s, "02_outlier_comparison_mundka.png", 6.4, 1.8, 6.65, 3.4)

# station boxplot below
add_image(s, "02_outlier_boxplot_stations.png", 6.4, 5.25, 6.65, 1.9)
add_text(s, "Per-station outlier spread (justifies station-level clipping)",
         6.4, 7.1, 6.65, 0.28, size=9, color=C_GRAY, align=PP_ALIGN.CENTER)

footer_bar(s); slide_number(s, 4)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — EDA: SEASONALITY, DIURNAL, SPATIAL
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
slide_header(s, "Exploratory Data Analysis — Key Patterns",
             "Seasonality · Diurnal rhythm · Spatial consistency")

# three panels
panel_w = 4.15
for i, (fname, title, caption) in enumerate([
    ("03_monthly_boxplots_all_stations.png",
     "Seasonal Pattern",
     "3× Winter–Summer PM2.5 ratio"),
    ("05_diurnal_by_season.png",
     "Diurnal Cycle",
     "Twin peaks: 8 AM rush + 9 PM boundary-layer collapse"),
    ("07_correlation_heatmap.png",
     "Spatial Correlation",
     "r > 0.85 across all stations → city-wide phenomenon"),
]):
    x = 0.2 + i * (panel_w + 0.15)
    add_rect(s, x, 1.35, panel_w, 5.55, fill=C_SURFACE,
             line=RGBColor(0x2A, 0x3B, 0x60))
    add_text(s, title, x+0.1, 1.42, panel_w-0.2, 0.42,
             size=13, bold=True, color=C_ACCENT2, align=PP_ALIGN.CENTER)
    add_image(s, fname, x+0.1, 1.88, panel_w-0.2, 3.8)
    add_text(s, caption, x+0.05, 5.7, panel_w-0.1, 0.5,
             size=10, color=C_LIGHT, align=PP_ALIGN.CENTER)

# bonus caption
add_text(s, "💡 Strong seasonality, diurnal rhythm & high spatial correlation → all three are directly exploited in feature engineering",
         0.3, 6.88, 12.7, 0.35, size=10, color=C_GRAY, align=PP_ALIGN.CENTER)

footer_bar(s); slide_number(s, 5)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — EDA: DIWALI SPIKE
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
slide_header(s, "Event Analysis — Diwali Spike 🎆",
             "Diwali causes 2–5× the winter seasonal baseline across all 4 years")

# hero image top
add_image(s, "06_diwali_event_windows.png", 0.25, 1.35, 8.5, 3.9)
add_text(s, "10-day PM2.5 windows centred on Diwali (2021–2024)",
         0.25, 5.2, 8.5, 0.3, size=9, color=C_GRAY, align=PP_ALIGN.CENTER)

# right insets
add_text(s, "Residual Analysis", 9.0, 1.35, 4.1, 0.38,
         size=12, bold=True, color=C_ACCENT2)
add_image(s, "10_diwali_event_analysis.png", 9.0, 1.78, 4.1, 2.25)
add_text(s, "Spike above de-seasonalised baseline", 9.0, 3.98, 4.1, 0.3,
         size=9, color=C_GRAY, align=PP_ALIGN.CENTER)

add_text(s, "Stubble Burning Season", 9.0, 4.45, 4.1, 0.38,
         size=12, bold=True, color=C_ACCENT2)
add_image(s, "06_stubble_burning_analysis.png", 9.0, 4.88, 4.1, 1.55)
add_text(s, "Oct–Nov co-occurrence with Diwali", 9.0, 6.38, 4.1, 0.28,
         size=9, color=C_GRAY, align=PP_ALIGN.CENTER)

# key insight box
add_rect(s, 0.25, 5.6, 8.5, 1.1, fill=RGBColor(0x1E, 0x14, 0x02),
         line=C_ACCENT)
add_text(s, "⚡  Motivates: days_since_diwali feature · is_stubble_season flag · Failure Mode CUSUM analysis",
         0.45, 5.75, 8.1, 0.75, size=13, color=C_ACCENT2)

footer_bar(s); slide_number(s, 6)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — FEATURE ENGINEERING
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
slide_header(s, "Feature Engineering",
             "63 domain-aware features bridging physics and ML")

feat_rows = [
    ("Lag Features",       "pm25_lag1, lag2, lag3, lag7, lag14, lag30",       "Persistence + weekly cycle"),
    ("Rolling Statistics", "3d/7d/14d/30d mean, std, max",                    "Multi-scale volatility"),
    ("Cyclical Encoding",  "sin/cos of month, hour, day-of-week",             "Prevents Dec–Jan discontinuity"),
    ("Wind Decomposition", "U = −ws·sin(θ)    V = −ws·cos(θ)",               "NW smoke transport from Punjab"),
    ("Event Flags",        "days_since_diwali (capped@10), is_stubble_season","Physics-driven event signal"),
    ("Interaction",        "humidity × temperature",                           "Hygroscopic PM2.5 growth"),
]
add_table_simple(s, ["Feature Type", "Examples", "Rationale"],
                 feat_rows, 0.25, 1.38, 7.0, 3.3, text_size=11)

add_text(s, "Total: 63 features", 0.25, 4.7, 3.5, 0.38,
         size=13, bold=True, color=C_ACCENT2)

# right plots
add_text(s, "Feature → Target Correlation", 7.45, 1.38, 5.65, 0.35,
         size=12, bold=True, color=C_ACCENT2)
add_image(s, "feat_correlation_with_target.png", 7.45, 1.78, 5.65, 2.6)

add_text(s, "PM2.5 vs Meteorological Variables", 7.45, 4.45, 5.65, 0.35,
         size=12, bold=True, color=C_ACCENT2)
add_image(s, "07_pm25_vs_meteo.png", 7.45, 4.85, 5.65, 2.35)

# ablation note
add_rect(s, 0.25, 5.1, 6.95, 0.85, fill=C_SURFACE, line=RGBColor(0x2A, 0x3B, 0x60))
add_text(s, "Ablation study (08_ablation_study.png) confirms each feature group reduces MAE — "
            "wind decomposition alone accounts for ~3 µg/m³ improvement",
         0.4, 5.18, 6.6, 0.65, size=11, color=C_LIGHT)

footer_bar(s); slide_number(s, 7)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — MODEL ZOO
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
slide_header(s, "Model Zoo — What We Trained",
             "From statistical baselines to state-of-the-art Transformers")

models = [
    ("ARIMA /\nSARIMA",  "Statistical",    C_GRAY,    "Seasonal baseline\nGaussian assumption"),
    ("XGBoost\n(Global)","Tree Boosting",  C_GREEN,   "Station as feature\nCross-station transfer"),
    ("LSTM",             "Recurrent",      C_YELLOW,  "Sequential memory\nGradient flow"),
    ("Informer",         "Transformer",    C_ACCENT2, "ProbSparse attn\nO(L log L) complexity"),
    ("PatchTST",         "Transformer",    C_ACCENT,  "Patch tokenisation\nChannel independence ⭐"),
]
bw = 2.3
for i, (name, mtype, col, desc) in enumerate(models):
    x = 0.25 + i * (bw + 0.1)
    add_rect(s, x, 1.38, bw, 3.5, fill=C_SURFACE, line=col)
    # top colour strip
    add_rect(s, x, 1.38, bw, 0.22, fill=col)
    add_text(s, name,  x+0.1, 1.65, bw-0.2, 0.75,
             size=15, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s, mtype, x+0.1, 2.4, bw-0.2, 0.38,
             size=11, color=C_GRAY, align=PP_ALIGN.CENTER)
    add_rect(s, x+0.2, 2.78, bw-0.4, 0.025, fill=RGBColor(0x2A, 0x3B, 0x60))
    add_text(s, desc,  x+0.1, 2.85, bw-0.2, 0.9,
             size=11, color=C_LIGHT, align=PP_ALIGN.CENTER)

# arrow connector text
add_text(s, "→   Increasing model complexity   →",
         0.25, 4.95, 12.8, 0.4, size=12, color=C_GRAY, align=PP_ALIGN.CENTER)

# walk-forward CV
add_text(s, "Walk-Forward Cross-Validation Setup", 0.25, 5.42, 6.0, 0.38,
         size=12, bold=True, color=C_ACCENT2)
add_image(s, "09_walkforward_cv.png", 0.25, 5.85, 6.0, 1.35)

# learning curve
add_text(s, "XGBoost Data-Scaling Learning Curve", 6.45, 5.42, 6.6, 0.38,
         size=12, bold=True, color=C_ACCENT2)
add_image(s, "09_learning_curve_xgb.png", 6.45, 5.85, 6.6, 1.35)

footer_bar(s); slide_number(s, 8)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — RESULTS LEADERBOARD
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
slide_header(s, "Results — The Leaderboard 🏆",
             "Held-out 2024–2025 test set · 4-year training")

# hourly table
add_text(s, "Hourly Forecasting", 0.25, 1.38, 5.8, 0.38,
         size=14, bold=True, color=C_ACCENT2)
add_table_simple(s, ["Model", "MAE (µg/m³)", "R²"],
                 [("XGBoost Global", "39.47", "0.652"),
                  ("LSTM",           "38.19", "0.725"),
                  ("Informer",       "32.48", "0.777"),
                  ("PatchTST ⭐",    "30.42", "0.821")],
                 0.25, 1.82, 5.8, 1.85, text_size=12)

# daily table
add_text(s, "Daily Forecasting", 0.25, 3.78, 5.8, 0.38,
         size=14, bold=True, color=C_ACCENT2)
add_table_simple(s, ["Model", "MAE (µg/m³)", "R²"],
                 [("Persistence",       "68.2",   "0.420"),
                  ("SARIMA",            "111.22", "−0.717"),
                  ("XGBoost Global ⭐", "33.62",  "0.570"),
                  ("LSTM ❌",           "89.67",  "−0.884"),
                  ("PatchTST",          "44.16",  "0.406")],
                 0.25, 4.22, 5.8, 2.0, text_size=12)

# main comparison chart
add_text(s, "Model Comparison (4-year training)", 6.25, 1.38, 6.85, 0.38,
         size=13, bold=True, color=C_ACCENT2)
add_image(s, "model_comparison_4yr.png", 6.25, 1.82, 6.85, 3.0)

# learning curves + per-station
add_image(s, "learning_curve_all_models.png", 6.25, 4.9, 3.3, 2.3)
add_text(s, "Data-scaling curves", 6.25, 7.15, 3.3, 0.25,
         size=9, color=C_GRAY, align=PP_ALIGN.CENTER)

add_image(s, "eval_per_station_mae.png", 9.75, 4.9, 3.35, 2.3)
add_text(s, "Per-station consistency", 9.75, 7.15, 3.35, 0.25,
         size=9, color=C_GRAY, align=PP_ALIGN.CENTER)

footer_bar(s); slide_number(s, 9)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — LSTM DAILY FAILURE
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
slide_header(s, "Deep Dive: Why Does LSTM Fail on Daily Data?",
             "An important negative result — sequence length matters")

add_bullet_box(s, [
    ("LSTM daily R² = −0.884  (worse than a Persistence baseline)", True, C_RED),
    "Only ~1,460 daily training steps over 4 years",
    "Recurrent gradient flow cannot learn seasonal structure at this scale",
    ("XGBoost overcomes this via explicit lag & calendar features", False, C_GREEN),
    "",
    ("Lesson:", True, C_ACCENT2),
    "Domain-aware feature engineering substitutes for recurrent inductive bias",
    "when sequence length is constrained — no architecture can fix data scarcity",
    "",
    ("Ensemble Stacking rescues the loss:", True, C_ACCENT2),
    ("R²: 0.543 → 0.683   |   MAE improvement: +5.75 µg/m³", False, C_ACCENT2),
], 0.3, 1.38, 5.8, 5.7, bullet="▸", size=14)

# residuals plot
add_text(s, "Residual Analysis (XGBoost vs LSTM)", 6.3, 1.38, 6.8, 0.38,
         size=13, bold=True, color=C_ACCENT2)
add_image(s, "eval_residuals.png", 6.3, 1.82, 6.8, 2.7)

# seasonal error
add_text(s, "Seasonal MAE Breakdown", 6.3, 4.6, 3.3, 0.35,
         size=12, bold=True, color=C_ACCENT2)
add_image(s, "14_seasonal_error.png", 6.3, 5.0, 3.3, 2.18)

# ensemble stacking
add_text(s, "Ensemble Stacking Result", 9.8, 4.6, 3.3, 0.35,
         size=12, bold=True, color=C_ACCENT2)
add_image(s, "13_ensemble_stacking.png", 9.8, 5.0, 3.3, 2.18)

footer_bar(s); slide_number(s, 10)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — CONFORMAL PREDICTION
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
slide_header(s, "Uncertainty Estimation — Conformal Prediction 🛡️",
             "Distribution-free coverage guarantees for public health decisions")

# left: flat
add_rect(s, 0.25, 1.38, 6.0, 5.75, fill=C_SURFACE, line=C_RED)
add_text(s, "❌  Flat Conformal (Baseline)", 0.4, 1.45, 5.7, 0.42,
         size=14, bold=True, color=C_RED, align=PP_ALIGN.CENTER)
add_image(s, "prediction_intervals.png", 0.35, 1.95, 5.8, 3.4)
add_bullet_box(s, [
    "Split conformal on 2023 calibration set (n=3,209)",
    "Distribution-free guarantee: ≥90% coverage",
    ("⚠  14.5 ppt under-coverage in Winter (heavy tails!)", True, C_RED),
    "Global quantile too narrow for extreme pollution events",
], 0.4, 5.42, 5.7, 1.55, bullet="▸", size=11)

# right: hierarchical
add_rect(s, 6.55, 1.38, 6.55, 5.75, fill=C_SURFACE,
         line=C_GREEN)
add_text(s, "✅  Hierarchical Seasonal Conformal (Novel)", 6.7, 1.45, 6.25, 0.42,
         size=14, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
add_image(s, "15_seasonal_conformal.png", 6.65, 1.95, 6.35, 3.4)
add_bullet_box(s, [
    "Separate quantile per season (Winter/Summer/Monsoon/Post-Monsoon)",
    ("✅  Overall coverage: 90.6%", True, C_GREEN),
    ("✅  Winter deficit fully corrected", True, C_GREEN),
    ("✅  11.2% interval width reduction", True, C_GREEN),
], 6.7, 5.42, 6.25, 1.55, bullet="▸", size=11)

footer_bar(s); slide_number(s, 11)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — FAILURE MODE ANALYSIS
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
slide_header(s, "Failure Mode Analysis — CUSUM Drift Detection ⚠️",
             "Knowing when the model fails is as important as knowing when it succeeds")

# hero image
add_image(s, "16_failure_mode_analysis.png", 0.25, 1.38, 8.4, 4.15)
add_text(s, "CUSUM statistic + dominant error regimes", 0.25, 5.48, 8.4, 0.28,
         size=9, color=C_GRAY, align=PP_ALIGN.CENTER)

# three regime cards (right)
regimes = [
    ("🎆 Diwali Proximity", "±7 days", "CUSUM fires 7 days pre-Diwali\nwithout any calendar label", C_ACCENT),
    ("💨 Calm Wind (<1 m/s)", "60%", "of top-5% failure events\nNo transport dilution", C_YELLOW),
    ("❄️ Winter Inversion", "98%", "of all failures in Winter\n3.28× Winter/Summer MAE", C_RED),
]
for i, (title, stat, desc, col) in enumerate(regimes):
    y = 1.38 + i * 1.85
    add_rect(s, 8.85, y, 4.25, 1.65, fill=C_SURFACE, line=col)
    add_text(s, title, 8.95, y+0.1, 3.0, 0.38, size=13, bold=True, color=col)
    add_text(s, stat,  12.1, y+0.1, 0.9, 0.55, size=22, bold=True,
             color=col, align=PP_ALIGN.RIGHT)
    add_text(s, desc,  8.95, y+0.5, 4.05, 0.85, size=11, color=C_LIGHT)

# bottom strip: masked MAPE + DM test
add_image(s, "14_seasonal_error.png", 0.25, 5.78, 4.1, 1.38)
add_image(s, "12_significance_tests.png", 4.5, 5.78, 4.15, 1.38)
add_rect(s, 8.85, 5.78, 4.25, 1.38, fill=C_SURFACE, line=C_ACCENT2)
add_text(s, "Masked MAPE = 39.7%\n(PM2.5 > 10 µg/m³ filter)\n\nDiebold-Mariano p = 0.0029\nXGBoost > Persistence statistically",
         8.95, 5.88, 4.05, 1.18, size=11, color=C_LIGHT)

footer_bar(s); slide_number(s, 12)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — SHAP EXPLAINABILITY
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
slide_header(s, "Explainability — SHAP Analysis 🔍",
             "Opening the black box: what drives Delhi's PM2.5 predictions?")

# beeswarm (main)
add_text(s, "SHAP Beeswarm — Daily XGBoost", 0.25, 1.38, 7.5, 0.38,
         size=13, bold=True, color=C_ACCENT2)
add_image(s, "shap_daily_summary.png", 0.25, 1.82, 7.5, 4.05)

# bar chart (top right)
add_text(s, "Mean |SHAP| Ranking", 8.0, 1.38, 5.1, 0.38,
         size=12, bold=True, color=C_ACCENT2)
add_image(s, "shap_daily_bar.png", 8.0, 1.82, 5.1, 2.1)

# what-if (bottom right)
add_text(s, "What-If SHAP Sweep", 8.0, 4.0, 5.1, 0.38,
         size=12, bold=True, color=C_ACCENT2)
add_image(s, "shap_whatif.png", 8.0, 4.44, 5.1, 1.65)

# key insights box (bottom left)
add_rect(s, 0.25, 5.95, 7.5, 1.18, fill=RGBColor(0x0A, 0x14, 0x28),
         line=C_ACCENT)
add_bullet_box(s, [
    ("pm25_lag1 dominates  — pollution begets pollution", False, C_ACCENT2),
    "month_cos / is_winter drive seasonality",
    "wind_u (NW direction) → Punjab stubble smoke transport",
    "days_since_diwali adds signal beyond pure seasonality",
], 0.4, 6.02, 7.2, 1.05, bullet="▸", size=11)

footer_bar(s); slide_number(s, 13)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — ANOMALY DETECTION
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
slide_header(s, "Unsupervised Anomaly Detection 🤖",
             "Autoencoder trained on clean summer data — zero calendar labels")

# architecture box
add_rect(s, 0.25, 1.38, 5.0, 1.12, fill=C_SURFACE, line=C_ACCENT2)
add_text(s, "Architecture:  15 → 32 → 16 → 4 → 16 → 32 → 15",
         0.4, 1.48, 4.7, 0.42, size=13, bold=True, color=C_ACCENT2)
add_text(s, "Trained on Mar–Sep (clean summer) as 'normal' regime\nAnomaly score E = ‖x − f_dec(f_enc(x))‖²",
         0.4, 1.88, 4.7, 0.52, size=11, color=C_LIGHT)

# hero timeline
add_text(s, "Reconstruction Error Timeline — All 4 Diwali Dates Auto-Detected",
         0.25, 2.6, 12.85, 0.38, size=13, bold=True, color=C_ACCENT2)
add_image(s, "unsup_autoencoder_anomalies.png", 0.25, 3.02, 12.85, 3.15)

# VAE vs AE inset
add_text(s, "VAE vs. AE Anomaly Scores", 5.5, 1.38, 4.25, 0.38,
         size=12, bold=True, color=C_ACCENT2)
add_image(s, "11_vae_vs_ae_anomaly.png", 5.5, 1.82, 4.25, 1.05)

# t-SNE inset
add_text(s, "Latent Space t-SNE", 9.95, 1.38, 3.15, 0.38,
         size=12, bold=True, color=C_ACCENT2)
add_image(s, "14_latent_tsne.png", 9.95, 1.82, 3.15, 1.05)

# result banner
add_rect(s, 0.25, 6.22, 12.85, 0.78, fill=RGBColor(0x06, 0x18, 0x06),
         line=C_GREEN)
add_text(s, "✅  4-for-4 Diwali detection (2021, 2022, 2023, 2024) + stubble season "
            "— without any calendar label. These episodes represent genuine physics shifts.",
         0.45, 6.32, 12.45, 0.58, size=13, color=C_GREEN)

footer_bar(s); slide_number(s, 14)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — STATION CLUSTERING
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
slide_header(s, "Station Clustering — Spatial Intelligence 🏙️",
             "K-Means k=3 · Silhouette Score 0.61 · Elbow validated")

# cluster plot
add_text(s, "K-Means Cluster Profiles (monthly PM2.5)", 0.25, 1.38, 7.5, 0.38,
         size=13, bold=True, color=C_ACCENT2)
add_image(s, "unsup_kmeans_clusters.png", 0.25, 1.82, 7.5, 3.5)

# station embedding
add_text(s, "Station Embedding (reduced-dim)", 8.0, 1.38, 5.1, 0.38,
         size=12, bold=True, color=C_ACCENT2)
add_image(s, "unsup_embedding.png", 8.0, 1.82, 5.1, 3.5)

# cluster cards
clusters = [
    ("🔴 Cluster 1 — High Pollution", "Anand Vihar · Mundka",
     "Transport hubs, >250 µg/m³ winter, slow clearance", C_RED),
    ("🟡 Cluster 2 — Moderate / Mixed Urban",
     "Ashok Vihar · Bawana · Jahangirpuri · Rohini · Wazirpur · Punjabi Bagh",
     "Mid-range winter levels", C_YELLOW),
    ("🟢 Cluster 3 — Low Pollution / Windward", "Dwarka-Sector 8",
     "Lowest baseline, fastest clearance under westerly winds", C_GREEN),
]
cw = 4.2
for i, (title, stn, desc, col) in enumerate(clusters):
    x = 0.25 + i * (cw + 0.1)
    add_rect(s, x, 5.43, cw, 1.72, fill=C_SURFACE, line=col)
    add_text(s, title, x+0.1, 5.5,  cw-0.2, 0.38, size=12, bold=True, color=col)
    add_text(s, stn,   x+0.1, 5.88, cw-0.2, 0.42, size=10, color=C_ACCENT2)
    add_text(s, desc,  x+0.1, 6.28, cw-0.2, 0.75, size=10, color=C_LIGHT)

footer_bar(s); slide_number(s, 15)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — DASHBOARD
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
slide_header(s, "Product — Delhi AQ Dashboard 🚀",
             "End-to-end deployment: Streamlit · HuggingFace · Docker")

# two mode cards
for i, (icon, mode, audience, feats, col) in enumerate([
    ("🔬", "Lab Results", "Researchers / Evaluators",
     ["Forecast plots (all models)", "SHAP explanations",
      "Conformal prediction intervals", "Station clustering"],
     C_ACCENT2),
    ("📱", "Try Yourself", "Public / Policymakers",
     ["Weather-app AQI health cards", "Animated AQI gauge",
      "7-day outlook simulation", "Recursive lag-shift engine"],
     C_GREEN),
]):
    x = 0.25 + i * 6.6
    add_rect(s, x, 1.38, 6.3, 4.4, fill=C_SURFACE, line=col)
    add_text(s, f"{icon}  {mode}", x+0.2, 1.48, 5.9, 0.55,
             size=18, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s, f"Audience: {audience}", x+0.2, 2.05, 5.9, 0.38,
             size=11, color=C_GRAY, align=PP_ALIGN.CENTER)
    add_rect(s, x+0.2, 2.45, 5.9, 0.02, fill=RGBColor(0x2A, 0x3B, 0x60))
    for j, feat in enumerate(feats):
        add_text(s, f"▸  {feat}", x+0.3, 2.55 + j*0.47, 5.7, 0.42,
                 size=12, color=C_LIGHT)

# tech stack
add_rect(s, 0.25, 5.88, 12.85, 0.82, fill=C_SURFACE, line=RGBColor(0x2A, 0x3B, 0x60))
add_text(s,
         "⚙️  Stack: Streamlit · Python · Scikit-learn · PyTorch · SHAP  |  "
         "🤗 Models on HuggingFace Hub  |  🐳 Docker (port 7860)",
         0.45, 5.98, 12.45, 0.58, size=12, color=C_LIGHT)

# URL
add_text(s, "🌐  huggingface.co/spaces/Guna-Venkat-Doddi-251140009/delhi-aq-dashboard",
         0.25, 6.78, 12.85, 0.42, size=13, bold=True, color=C_ACCENT,
         align=PP_ALIGN.CENTER)

footer_bar(s); slide_number(s, 16)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — KEY CONTRIBUTIONS
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
slide_header(s, "Key Contributions — What Is New in This Work",
             "7 novel contributions beyond standard coursework")

contribs = [
    ("1", "Hierarchical Seasonal Conformal Prediction",
     "Corrects 14.5 ppt Winter deficit · 11.2% tighter intervals", C_ACCENT2),
    ("2", "CUSUM Failure Mode Analysis",
     "3 dominant error regimes · actionable deployment guardrails", C_ACCENT),
    ("3", "Masked MAPE (PM2.5>10 filter)",
     "Stable cross-model metric · avoids summer near-zero inflation", C_LIGHT),
    ("4", "Global XGBoost > per-station",
     "+1–2 MAE units · confirmed cross-station transfer learning", C_GREEN),
    ("5", "LSTM daily failure (negative result)",
     "R²=−0.884 → sequence-length threshold discovered", C_RED),
    ("6", "Unsupervised Diwali detection",
     "4-for-4 hit rate (2021–2024) · zero calendar labels", C_YELLOW),
    ("7", "End-to-end live dashboard",
     "Consumer + Research modes · Docker · HuggingFace deployed", C_ACCENT2),
]
rh = 0.72
for i, (num, title, impact, col) in enumerate(contribs):
    y = 1.42 + i * (rh + 0.06)
    add_rect(s, 0.25, y, 0.55, rh, fill=col)
    add_text(s, num, 0.25, y + 0.12, 0.55, rh-0.24,
             size=18, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    add_rect(s, 0.85, y, 12.25, rh, fill=C_SURFACE, line=RGBColor(0x2A, 0x3B, 0x60))
    add_text(s, title, 1.0, y+0.06, 6.0, 0.38, size=13, bold=True, color=col)
    add_text(s, impact, 7.2, y+0.06, 5.7, 0.38, size=12, color=C_LIGHT,
             align=PP_ALIGN.RIGHT)

footer_bar(s); slide_number(s, 17)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — CONCLUSION
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
slide_header(s, "Conclusion & Key Takeaways",
             "Bridging ML research and actionable environmental monitoring")

takeaways = [
    ("🏆 Best Hourly Model:", "PatchTST  —  MAE 30.42 µg/m³, R² = 0.821",  C_GREEN),
    ("🏆 Best Daily Model:",  "XGBoost (Global)  —  MAE 33.62 µg/m³, R² = 0.570", C_GREEN),
    ("⚠️  Negative Result:",  "LSTM daily R² = −0.884  →  sequence length determines architecture choice", C_RED),
    ("🛡️  Calibrated Uncertainty:", "Seasonal conformal prediction achieves 90.6% empirical coverage", C_ACCENT2),
    ("🔍  Failure Modes:", "Physics-driven: calm wind + winter inversion = dominant error regime", C_YELLOW),
    ("🤖  Unsupervised:", "Autoencoder autonomously identifies all 4 Diwali events without labels", C_ACCENT2),
    ("🌐  Deployed:", "Live dashboard — dual Consumer + Research modes on HuggingFace", C_GREEN),
]
y = 1.42
for label, detail, col in takeaways:
    add_rect(s, 0.25, y, 12.85, 0.68, fill=C_SURFACE, line=RGBColor(0x2A, 0x3B, 0x60))
    add_text(s, label,  0.4, y+0.12, 3.2, 0.45, size=13, bold=True, color=col)
    add_text(s, detail, 3.65, y+0.12, 9.3, 0.45, size=13, color=C_WHITE)
    y += 0.76

# closing quote
add_rect(s, 0.25, 6.88, 12.85, 0.88, fill=RGBColor(0x15, 0x0C, 0x01),
         line=C_ACCENT)
add_text(s,
         '"We don\'t just predict pollution — we know when our model will fail, '
         'why it will fail, and what to do about it."',
         0.6, 6.96, 12.25, 0.72, size=14, bold=True, color=C_ACCENT,
         align=PP_ALIGN.CENTER)

footer_bar(s); slide_number(s, 18)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Q&A / THANK YOU
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)

# faded background
add_rect(s, 0, 0, 0.5, 7.5, fill=C_ACCENT)

img_path = PLOTS / "14_latent_tsne.png"
if img_path.exists():
    pic = s.shapes.add_picture(str(img_path), Inches(0), Inches(0),
                                Inches(13.33), Inches(7.5))
    sp_tree = s.shapes._spTree
    sp_tree.remove(pic._element)
    sp_tree.insert(2, pic._element)

add_rect(s, 0, 0, 13.33, 7.5, fill=RGBColor(0x0D, 0x12, 0x1F))
# re-add accent bar on top of overlay
add_rect(s, 0, 0, 0.5, 7.5, fill=C_ACCENT)

add_text(s, "Thank You", 0.7, 1.5, 12.0, 1.1,
         size=52, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Questions & Discussion", 0.7, 2.65, 12.0, 0.65,
         size=24, color=C_ACCENT2, align=PP_ALIGN.CENTER)

add_rect(s, 2.5, 3.45, 8.33, 0.05, fill=C_ACCENT)

add_text(s, "Guna Venkat Doddi  ·  Roll No. 251140009", 0.7, 3.65, 12.0, 0.48,
         size=16, color=C_LIGHT, align=PP_ALIGN.CENTER)
add_text(s, "DMS673: Applied Machine Learning  ·  IIT Kanpur  ·  April 2026",
         0.7, 4.12, 12.0, 0.42, size=14, color=C_GRAY, align=PP_ALIGN.CENTER)

add_text(s, "🌐  Live Demo:", 2.0, 4.85, 2.5, 0.45,
         size=14, color=C_ACCENT2, align=PP_ALIGN.RIGHT)
add_text(s, "huggingface.co/spaces/Guna-Venkat-Doddi-251140009/delhi-aq-dashboard",
         4.6, 4.85, 8.5, 0.45, size=14, bold=True, color=C_ACCENT)

add_text(s, "AI Assistance: Claude · DeepSeek · Gemini (used for ideation, code review & report writing)",
         0.7, 5.55, 12.0, 0.38, size=11, color=C_GRAY, align=PP_ALIGN.CENTER)

slide_number(s, 19)

# ─── Save ────────────────────────────────────────────────────────────────────
prs.save(str(OUT))
print(f"\n✅  Saved: {OUT}")
print(f"   Slides: 19")
print(f"   Size:   {OUT.stat().st_size / 1024:.1f} KB")
